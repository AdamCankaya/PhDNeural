"""Generate FAU-branded Bio-NAS thesis committee intro slides."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# FAU Visual Standards (primary identity colors)
FAU_BLUE = RGBColor(0x00, 0x33, 0x66)
FAU_RED = RGBColor(0xCC, 0x00, 0x00)
FAU_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
FAU_SILVER = RGBColor(0xA7, 0xA9, 0xAC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE = RGBColor(0xF7, 0xF8, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x55, 0x66)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT = Path(__file__).with_name("Bio-NAS_Thesis_Committee_Intro.pptx")


def _set_run_font(run, *, size_pt: float, bold: bool = False, color: RGBColor = DARK_TEXT, name: str = "Calibri"):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.get_or_add_bodyPr().set(qn("a:anchor"), {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=size, bold=bold, color=color, name=font)
    return box


def _add_bullets(slide, left, top, width, height, items, *, size=16, color=DARK_TEXT, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(10)
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"•  {item}"
        _set_run_font(run, size_pt=size, bold=(bold_first and i == 0), color=color)
    return box


def _chrome(slide, title: str, subtitle: str | None = None, *, dark_title=False):
    """FAU header bar + red accent + footer."""
    # Top blue bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FAU_BLUE
    bar.line.fill.background()

    # Red accent under header
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SLIDE_W, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = FAU_RED
    accent.line.fill.background()

    # Left red rail
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.03), Inches(0.12), SLIDE_H - Inches(1.03))
    rail.fill.solid()
    rail.fill.fore_color.rgb = FAU_RED
    rail.line.fill.background()

    # Footer
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.42), SLIDE_W, Inches(0.42))
    foot.fill.solid()
    foot.fill.fore_color.rgb = FAU_BLUE
    foot.line.fill.background()

    _add_textbox(
        slide,
        Inches(0.35),
        Inches(0.18),
        Inches(12.5),
        Inches(0.55),
        title,
        size=26,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        _add_textbox(
            slide,
            Inches(0.35),
            Inches(0.58),
            Inches(12.5),
            Inches(0.32),
            subtitle,
            size=12,
            color=FAU_GRAY,
        )

    _add_textbox(
        slide,
        Inches(0.35),
        SLIDE_H - Inches(0.38),
        Inches(8.5),
        Inches(0.32),
        "Adam Cankaya  ·  Florida Atlantic University  ·  College of Engineering and Computer Science",
        size=10,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        Inches(9.2),
        SLIDE_H - Inches(0.38),
        Inches(3.8),
        Inches(0.32),
        "acankaya2017@fau.edu",
        size=10,
        color=FAU_GRAY,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _set_notes(slide, script: str):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = script.strip()
    run.font.size = Pt(12)
    run.font.name = "Calibri"


def _bg(slide, color=NEAR_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _card(slide, left, top, width, height, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = FAU_GRAY
    shape.line.width = Pt(1)
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ----- Slide 1: Title -----
    s = prs.slides.add_slide(blank)
    _bg(s, FAU_BLUE)
    # Red bottom band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(1.55), SLIDE_W, Inches(1.55))
    band.fill.solid()
    band.fill.fore_color.rgb = FAU_RED
    band.line.fill.background()
    # Thin silver line
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(1.55), SLIDE_W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = FAU_GRAY
    line.line.fill.background()

    _add_textbox(s, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.4), "FLORIDA ATLANTIC UNIVERSITY", size=14, bold=True, color=FAU_GRAY)
    _add_textbox(s, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.35), "College of Engineering and Computer Science", size=16, color=WHITE)
    _add_textbox(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(1.6),
        "Spatio-Temporal Bio-NAS for\nMulti-Omic Disease Prediction",
        size=36,
        bold=True,
        color=WHITE,
    )
    _add_textbox(
        s,
        Inches(0.8),
        Inches(3.9),
        Inches(11.5),
        Inches(0.7),
        "Dual-Track Neural Architecture Search Constrained by Biological Pathways",
        size=18,
        color=FAU_GRAY,
    )
    _add_textbox(
        s,
        Inches(0.8),
        SLIDE_H - Inches(1.25),
        Inches(11.5),
        Inches(0.9),
        "Adam Cankaya\nacankaya2017@fau.edu  ·  Thesis Committee Introduction  ·  Fall 2026",
        size=16,
        bold=True,
        color=WHITE,
    )
    _set_notes(
        s,
        """Good afternoon, and thank you for joining. My name is Adam Cankaya. I am a doctoral student in the College of Engineering and Computer Science at Florida Atlantic University. You can reach me at acankaya2017@fau.edu.

Today I want to introduce my dissertation project, Spatio-Temporal Bio-NAS for Multi-Omic Disease Prediction. In short, this work asks whether building biological knowledge into neural architecture search improves how we model and forecast disease from multi-omic data over time.

I will cover the research question, the dual-track experimental design, the five disease categories, the technical pipeline, the three-year roadmap, deliverables, current status, and a few specific asks for the committee. The goal of this session is orientation and feedback, not a full defense of results.""",
    )

    # ----- Slide 2: Research question -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Research Question & Hypothesis")
    _card(s, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.35))
    _add_textbox(
        s,
        Inches(0.75),
        Inches(1.55),
        Inches(11.8),
        Inches(1.0),
        "Does biological etiology dictate optimal neural architecture?",
        size=26,
        bold=True,
        color=FAU_BLUE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    items = [
        "Hypothesis: Bio-NAS (KEGG / Reactome pathway constraints) outperforms unconstrained NAS on accuracy, interpretability, and/or sparsity.",
        "Setting: spatio-temporal deep learning over irregular longitudinal multi-omic visits.",
        "Strategy: prove the method on BRCA first, then scale to four additional disease categories if justified.",
    ]
    _add_bullets(s, Inches(0.7), Inches(3.0), Inches(11.8), Inches(3.2), items, size=18)
    _set_notes(
        s,
        """The central research question is: does biological etiology dictate optimal neural architecture? In other words, when we search for neural networks that predict disease from multi-omic data, should biology shape that search, or is unconstrained mathematical optimization enough?

My working hypothesis is that Bio-NAS—constraining artificial synapses with known pathway structure from resources like KEGG and Reactome—will outperform unconstrained neural architecture search. I will look for gains in predictive accuracy, interpretability, and computational sparsity. Any one of those advantages can matter scientifically; all three would be ideal.

The modeling setting is deliberately longitudinal and spatio-temporal. Molecular measurements arrive at irregular intervals, so time itself is part of the architecture problem. Operationally, I validate the full dual-track method on breast cancer first, then expand only if the BRCA evidence supports it.""",
    )

    # ----- Slide 3: Dual tracks -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Dual-Track Experimental Design", "Track A = control  ·  Track B = innovation  ·  Shared infrastructure for both")

    # Track A card
    _card(s, Inches(0.5), Inches(1.4), Inches(5.9), Inches(4.6))
    header_a = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(5.9), Inches(0.7))
    header_a.fill.solid()
    header_a.fill.fore_color.rgb = FAU_BLUE
    header_a.line.fill.background()
    _add_textbox(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.5), "Track A — Standard NAS", size=20, bold=True, color=WHITE)
    _add_bullets(
        s,
        Inches(0.75),
        Inches(2.3),
        Inches(5.4),
        Inches(3.4),
        [
            "Role: control / baseline",
            "Unconstrained search over layers, width, dropout",
            "Optimizes for mathematical performance",
            "Establishes what Optuna finds without biology",
        ],
        size=16,
    )

    # Track B card
    _card(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(4.6))
    header_b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.7))
    header_b.fill.solid()
    header_b.fill.fore_color.rgb = FAU_RED
    header_b.line.fill.background()
    _add_textbox(s, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5), "Track B — Bio-NAS", size=20, bold=True, color=WHITE)
    _add_bullets(
        s,
        Inches(7.15),
        Inches(2.3),
        Inches(5.4),
        Inches(3.4),
        [
            "Role: biologically informed innovation",
            "Pathway adjacency → MaskedLinear layers",
            "Forces search toward true biological routes",
            "Tests gains in accuracy, interpretability, sparsity",
        ],
        size=16,
    )
    _set_notes(
        s,
        """This project is structured as a rigid dual-track A/B test. Track A is standard neural architecture search: Optuna explores an unconstrained space of layers, widths, and regularization, optimizing purely for predictive performance. That track is the control.

Track B is Bio-NAS. Here I translate biological blueprints—primarily KEGG and Reactome pathways—into adjacency masks and apply them through MaskedLinear layers. The search is still Optuna-driven, but the candidate networks are forced to respect biological connectivity.

Both tracks share the same data splits, tensors, and compute stack, so differences can be attributed to the biological constraints rather than preprocessing artifacts. The scientific claim is comparative: does Track B beat Track A on holdout metrics that matter clinically and scientifically?""",
    )

    # ----- Slide 4: Diseases -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Disease Comparative Matrix", "Five categories  ·  BRCA is the Year 1–2 anchor  ·  Other 4 gated by BRCA evidence")

    rows = [
        ("Oncological (Anchor)", "Breast Invasive Carcinoma (BRCA)", "TCGA / AURORA"),
        ("Neurological", "Alzheimer’s Disease", "ADNI"),
        ("Autoimmune", "Rheumatoid Arthritis", "GSE138747"),
        ("Metabolic", "Type 2 Diabetes", "KORA F4/FF4"),
        ("Aging / Epigenetic", "Epigenetic Aging", "LBC1936"),
    ]
    # Table header
    y0 = Inches(1.35)
    col_x = [Inches(0.5), Inches(3.6), Inches(8.4)]
    col_w = [Inches(3.0), Inches(4.7), Inches(4.4)]
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y0, Inches(12.3), Inches(0.5))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = FAU_BLUE
    hdr.line.fill.background()
    for x, w, label in zip(col_x, col_w, ["Category", "Disease", "Primary cohort source"]):
        _add_textbox(s, x + Inches(0.1), y0 + Inches(0.08), w - Inches(0.15), Inches(0.35), label, size=14, bold=True, color=WHITE)

    for i, (cat, dis, src) in enumerate(rows):
        top = y0 + Inches(0.5) + Inches(i * 0.72)
        bg = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xF2, 0xF7)
        row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(12.3), Inches(0.72))
        row.fill.solid()
        row.fill.fore_color.rgb = bg
        row.line.color.rgb = FAU_GRAY
        row.line.width = Pt(0.75)
        vals = [cat, dis, src]
        colors = [FAU_RED if i == 0 else FAU_BLUE, DARK_TEXT, MUTED]
        bolds = [True, True, False]
        for x, w, val, col, b in zip(col_x, col_w, vals, colors, bolds):
            _add_textbox(s, x + Inches(0.1), top + Inches(0.18), w - Inches(0.15), Inches(0.4), val, size=14, bold=b, color=col)

    _set_notes(
        s,
        """The comparative matrix spans five disease categories chosen for biological and tempo diversity. Breast invasive carcinoma is the anchor. Years one and two complete a full dual-track vertical slice on BRCA using TCGA for multi-omic richness and AURORA when we need true paired molecular timepoints.

The other four are Alzheimer’s via ADNI, rheumatoid arthritis via an open multi-omic GEO series, type 2 diabetes via KORA, and epigenetic aging via the Lothian Birth Cohort. These represent neurological, autoimmune, metabolic, and aging contexts.

Importantly, dual-track work on the other four is gated. Only if BRCA holdout shows a Track B advantage in accuracy, interpretability, or sparsity do we justify expanding the comparative claims. A null BRCA result is still a thesis contribution; it simply narrows multi-disease claims.""",
    )

    # ----- Slide 5: Technical approach -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Technical Approach", "From longitudinal cohorts to searchable spatio-temporal architectures")

    steps = [
        ("1", "Cohorts &\nholdout", "Patient-level 80/20\nno leakage"),
        ("2", "Feature\nmaps", "Spatial sites ×\nvisits & Δt"),
        ("3", "4D HDF5\ntensors", "(B, T, S, C)\nmulti-omic"),
        ("4", "NAS\nmodules", "CNN / Transformer\n+ ConvLSTM / Attn"),
        ("5", "Optuna &\ncausal CV", "Track A then B\ndistributed search"),
        ("6", "Eval &\nattr.", "Holdout A vs B\nCaptum / Streamlit"),
    ]
    for i, (num, title, body) in enumerate(steps):
        left = Inches(0.4 + i * 2.15)
        _card(s, left, Inches(1.55), Inches(2.0), Inches(4.3))
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.65), Inches(1.75), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = FAU_RED if i % 2 else FAU_BLUE
        circle.line.fill.background()
        _add_textbox(s, left + Inches(0.65), Inches(1.85), Inches(0.7), Inches(0.5), num, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_textbox(s, left + Inches(0.1), Inches(2.65), Inches(1.8), Inches(1.1), title, size=15, bold=True, color=FAU_BLUE, align=PP_ALIGN.CENTER)
        _add_textbox(s, left + Inches(0.1), Inches(3.85), Inches(1.8), Inches(1.5), body, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(1.95), Inches(3.4), Inches(0.22), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = FAU_RED
            arrow.line.fill.background()

    _set_notes(
        s,
        """Methodologically, the pipeline has six stages. First, we inventory longitudinal multi-omic cohorts and lock a patient-level holdout before any fitting, so there is no subject leakage and no peeking at future visits.

Second and third, we map spatial axes—such as CpG sites or genomic coordinates—and temporal axes—visit index and inter-visit delta-t—into four-dimensional tensors stored in HDF5 with shape batch, time, spatial features, and channels.

Fourth, the NAS search space includes spatial modules like one-dimensional CNNs and spatial transformers, plus temporal modules such as ConvLSTM, GRU, and temporal attention that consume delta-t embeddings. Fifth, Optuna runs Track A then Track B with causal cross-validation on distributed workers. Sixth, we evaluate on frozen holdout trajectories, run ablations and classical baselines, then attribute predictions with Captum or SHAP and surface trajectories in a Streamlit demo.""",
    )

    # ----- Slide 6: Roadmap -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Three-Year Roadmap", "Fall 2026 → Summer 2029  ·  24 checklist tasks synced to GitHub issues")

    years = [
        ("Year 1", "Fall 2026 – Summer 2027", FAU_BLUE, [
            "All-5 cohort inventory",
            "BRCA feature maps & spacing",
            "Leakage-safe 4D ETL + Δt",
            "Compute stack & Optuna hub",
        ], "Shared (A+B) · BRCA-first"),
        ("Year 2", "Fall 2027 – Summer 2028", FAU_RED, [
            "Spatial & temporal modules",
            "Dual-track Optuna studies",
            "Holdout eval + ablations",
            "Classical baselines · scaling gate",
        ], "Track A vs B · BRCA"),
        ("Year 3", "Fall 2028 – Summer 2029", FAU_BLUE, [
            "Attribution & driver maps",
            "Trajectory dashboard",
            "Taxonomy / tempo compare",
            "Dissertation + OSS release",
        ], "BRCA → Other 4 / All 5"),
    ]
    for i, (yr, dates, color, bullets, scope) in enumerate(years):
        left = Inches(0.45 + i * 4.25)
        _card(s, left, Inches(1.35), Inches(4.05), Inches(4.85))
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.35), Inches(4.05), Inches(1.05))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        _add_textbox(s, left + Inches(0.2), Inches(1.42), Inches(3.65), Inches(0.45), yr, size=22, bold=True, color=WHITE)
        _add_textbox(s, left + Inches(0.2), Inches(1.9), Inches(3.65), Inches(0.35), dates, size=12, color=FAU_GRAY)
        _add_bullets(s, left + Inches(0.25), Inches(2.6), Inches(3.55), Inches(2.6), bullets, size=14)
        _add_textbox(s, left + Inches(0.2), Inches(5.55), Inches(3.65), Inches(0.45), scope, size=11, bold=True, color=color)

    _set_notes(
        s,
        """The program runs on a three-year academic calendar from Fall 2026 through Summer 2029, with twenty-four roadmap tasks tracked one-to-one as GitHub issues.

Year one is shared infrastructure: inventory all five diseases, build BRCA-first feature maps, construct leakage-safe four-dimensional tensors with delta-t embeddings, and stand up the compute and Optuna orchestration hub.

Year two engineers the search space—spatial then temporal modules—runs Track A and Track B Optuna studies on BRCA, and closes with holdout evaluation, ablations, and classical baselines. That summer comparison is the scaling gate.

Year three focuses on interpretability, a trajectory dashboard, structural taxonomy across disease tempos, the dissertation, and an open-source framework release. Other-four dual-track expansion happens only after the BRCA gate.""",
    )

    # ----- Slide 7: Goals -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Goals & Deliverables")

    goals = [
        ("Spatio-temporal NAS", "Optuna-driven search over spatial and temporal PyTorch modules"),
        ("Causal evaluation", "Patient holdout, causal CV, irregular Δt embeddings"),
        ("Comparative A/B evidence", "BRCA Track A vs B; Other 4 only if gate passes"),
        ("Interpretability", "CpG × time attribution maps and trajectory UI"),
        ("Open science", "Tagged Python OSS release plus dissertation"),
    ]
    for i, (title, desc) in enumerate(goals):
        top = Inches(1.3 + i * 0.9)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(0.18), Inches(0.75))
        bar.fill.solid()
        bar.fill.fore_color.rgb = FAU_RED if i % 2 else FAU_BLUE
        bar.line.fill.background()
        _card(s, Inches(0.7), top, Inches(12.1), Inches(0.75))
        _add_textbox(s, Inches(0.95), top + Inches(0.08), Inches(11.6), Inches(0.32), title, size=16, bold=True, color=FAU_BLUE)
        _add_textbox(s, Inches(0.95), top + Inches(0.38), Inches(11.6), Inches(0.3), desc, size=13, color=MUTED)

    _set_notes(
        s,
        """There are five concrete deliverables I want the committee to hold me to. First, a spatio-temporal NAS framework in PyTorch, orchestrated by Optuna, that can discover architectures for longitudinal multi-omic prediction.

Second, a causal evaluation protocol: patient-level holdout locked early, causal cross-validation, and explicit handling of irregular time through delta-t embeddings. Third, comparative A/B evidence—at minimum on BRCA—reporting whether Bio-NAS helps, and by how much, against unconstrained NAS and classical longitudinal baselines.

Fourth, interpretability artifacts: attribution maps that highlight which sites and timepoints drive forecasts, plus a Streamlit trajectory interface for demonstration. Fifth, an open-source release and the dissertation itself. I also want to be explicit that a negative or null Track B result on BRCA is still publishable science; it would tighten, not erase, the thesis contribution.""",
    )

    # ----- Slide 8: Status -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Current Status & Near-Term Next Steps")

    _card(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(4.7))
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(6.0), Inches(0.6))
    head.fill.solid()
    head.fill.fore_color.rgb = FAU_BLUE
    head.line.fill.background()
    _add_textbox(s, Inches(0.7), Inches(1.45), Inches(5.6), Inches(0.4), "Completed / in place", size=18, bold=True, color=WHITE)
    _add_bullets(
        s,
        Inches(0.75),
        Inches(2.2),
        Inches(5.5),
        Inches(3.5),
        [
            "Master plan + 24 GitHub-synced roadmap issues",
            "Live timeline dashboard and wiki",
            "Multi-disease cohort inventory drafted",
            "Code scaffolds: BRCA dataset, static MTL baseline, disease registry",
        ],
        size=15,
    )

    _card(s, Inches(6.85), Inches(1.35), Inches(6.0), Inches(4.7))
    head2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(1.35), Inches(6.0), Inches(0.6))
    head2.fill.solid()
    head2.fill.fore_color.rgb = FAU_RED
    head2.line.fill.background()
    _add_textbox(s, Inches(7.05), Inches(1.45), Inches(5.6), Inches(0.4), "Immediate next steps", size=18, bold=True, color=WHITE)
    _add_bullets(
        s,
        Inches(7.1),
        Inches(2.2),
        Inches(5.5),
        Inches(3.5),
        [
            "Lock primary cohorts and data-access accounts",
            "Finalize BRCA longitudinal matching strategy",
            "Build BRCA feature maps and genomic spacing",
            "Stand up leakage-safe ETL and Δt embeddings",
        ],
        size=15,
    )
    _set_notes(
        s,
        """On status: the project is past pure ideation and into an executable planning phase. The master plan, twenty-four roadmap issues, live dashboard, and wiki are in place. I have drafted a multi-disease cohort inventory with recommended primaries for BRCA, Alzheimer’s, RA, T2D, and epigenetic aging. Early code scaffolds exist for BRCA data loading, a static multi-task baseline, and a disease registry.

What comes next is operational. I need to lock cohorts and initiate the required data-access agreements, finalize how we match longitudinal samples for BRCA—especially given TCGA’s sparse molecular repeats versus AURORA’s paired primary–metastasis design—then build feature maps, genomic spacing rules, and the leakage-safe ETL path with delta-t embeddings. That is the Year-one vertical-slice foundation both tracks will consume.""",
    )

    # ----- Slide 9: Asks -----
    s = prs.slides.add_slide(blank)
    _bg(s)
    _chrome(s, "Asks for the Committee", "Feedback that will shape Year-1 priorities")

    asks = [
        "Does the dual-track framing and BRCA scaling gate feel appropriately rigorous for a dissertation contribution?",
        "How should we weight TCGA’s multi-omic richness against AURORA’s stronger molecular timepoints for the BRCA longitudinal design?",
        "Among ADNI, KORA, and LBC access paths, which Other-4 priorities should I pursue first while BRCA is underway?",
        "If Track B shows only interpretability or sparsity gains—not accuracy—should that still unlock multi-disease claims?",
    ]
    for i, ask in enumerate(asks):
        top = Inches(1.3 + i * 1.15)
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), top + Inches(0.15), Inches(0.55), Inches(0.55))
        num.fill.solid()
        num.fill.fore_color.rgb = FAU_RED if i % 2 else FAU_BLUE
        num.line.fill.background()
        _add_textbox(s, Inches(0.55), top + Inches(0.22), Inches(0.55), Inches(0.4), str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _card(s, Inches(1.3), top, Inches(11.5), Inches(1.0))
        _add_textbox(s, Inches(1.55), top + Inches(0.22), Inches(11.0), Inches(0.6), ask, size=15, color=DARK_TEXT)

    _set_notes(
        s,
        """I will close with four asks. First, I would value your judgment on whether the dual-track design and BRCA scaling gate are appropriately rigorous—strict enough to avoid overclaiming, but flexible enough for a complete dissertation if Track B is null.

Second, BRCA has a longitudinal tension: TCGA is excellent for multi-omic breadth but weak on repeated molecular sampling; AURORA is stronger on primary–metastasis pairs but smaller and more controlled-access. I would appreciate guidance on how to weight that tradeoff.

Third, while BRCA is underway, which Other-four access path should I prioritize—ADNI, KORA, or Lothian—so Year-three scaling is realistic. Fourth, if Bio-NAS improves interpretability or sparsity but not accuracy, should that still justify multi-disease expansion, or should accuracy remain the hard gate?

Thank you again. I welcome your questions and advice, and I’m happy to follow up at acankaya2017@fau.edu.""",
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
